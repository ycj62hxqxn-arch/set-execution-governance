"""
# ============================================================
# AI-AGENT OPERATIONS TOOL
# Tool ID    : GEN-PPTX-001
# Type       : Presentation / PDF Export Engine
# Actor      : YAI_AGENT (proposal / generation — no execution rights)
# Authority  : YAI may invoke this tool; output requires AA review before distribution
# ============================================================
# PURPOSE
#   Generates a fully-styled 16:9 PowerPoint (.pptx) presentation from
#   the BPB Internal Governance Architecture Report source document.
#   The output can also be exported to PDF via Keynote / PowerPoint "Export As".
#
# USAGE (AI Agent)
#   To regenerate after content changes:
#       "/Users/alaaatia/1. Project AI Manifest/.venv/bin/python" generate_pptx.py
#
# OUTPUT
#   BPB_Governance_Architecture_Report.pptx  (14 slides, 16:9 widescreen)
#
# DEPENDENCIES
#   python-pptx, Pillow  (installed in .venv)
#
# VERSION HISTORY
#   v1.0  2026-03-08  Initial generation — 14 slides, BPB dark-gold palette
# ============================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Palette ──────────────────────────────────────────────────────
BG      = RGBColor(0x0e, 0x0e, 0x1a)
DEEP    = RGBColor(0x1a, 0x1a, 0x2e)
MID     = RGBColor(0x2a, 0x2a, 0x3e)
GOLD    = RGBColor(0xc9, 0xa9, 0x61)
WHITE   = RGBColor(0xff, 0xff, 0xff)
MUTED   = RGBColor(0x88, 0x99, 0xaa)
GREEN   = RGBColor(0x27, 0xae, 0x60)
RED     = RGBColor(0xc0, 0x39, 0x2b)
BLUE    = RGBColor(0x34, 0x98, 0xdb)

W = Inches(13.33)   # 16:9 widescreen
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(layout)
    # dark background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    return txBox

def gold_bar(slide, top=Inches(0.08)):
    """Gold accent bar at top of slide."""
    add_rect(slide, 0, top, W, Inches(0.055), fill_color=GOLD)

def slide_label(slide, label_text, top=Inches(0.22)):
    add_text(slide, label_text.upper(), Inches(0.6), top, W - Inches(1.2), Inches(0.3),
             font_size=Pt(9), bold=True, color=GOLD, align=PP_ALIGN.LEFT)

def heading(slide, text, top=Inches(0.62), color=WHITE, size=Pt(32)):
    add_text(slide, text, Inches(0.6), top, W - Inches(1.2), Inches(0.7),
             font_size=size, bold=True, color=color, align=PP_ALIGN.LEFT)

def body_text(slide, text, left, top, width, height, size=Pt(13), color=MUTED):
    add_text(slide, text, left, top, width, height,
             font_size=size, bold=False, color=color, align=PP_ALIGN.LEFT)

def divider(slide, top=Inches(1.42)):
    add_rect(slide, Inches(0.6), top, Inches(0.55), Inches(0.045), fill_color=GOLD)

def footer_bar(slide):
    add_rect(slide, 0, H - Inches(0.42), W, Inches(0.42), fill_color=DEEP)
    add_text(slide, "BPB Solutions LTD  ·  Internal Governance Architecture Report  ·  CANONICAL_v1.0  ·  8 March 2026",
             Inches(0.3), H - Inches(0.36), W - Inches(0.6), Inches(0.3),
             font_size=Pt(8), color=GOLD, align=PP_ALIGN.CENTER)

# ── Slide builders ────────────────────────────────────────────────

def s1_title(prs):
    slide = blank_slide(prs)
    # accent rect left
    add_rect(slide, 0, 0, Inches(0.12), H, fill_color=GOLD)
    gold_bar(slide, top=Inches(0))

    # classification
    add_text(slide, "INTERNAL — CONFIDENTIAL", Inches(0.4), Inches(0.3), Inches(4), Inches(0.3),
             font_size=Pt(8), bold=True, color=GOLD)

    # title
    add_text(slide, "BPB Solutions LTD", Inches(0.6), Inches(1.4), W - Inches(1), Inches(1.2),
             font_size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "Internal Governance Architecture Report", Inches(0.6), Inches(2.7), W - Inches(1), Inches(0.55),
             font_size=Pt(20), bold=False, color=GOLD, align=PP_ALIGN.LEFT)

    add_text(slide, "Comprehensive Analysis  ·  Last 60 Days  ·  CANONICAL_v1.0",
             Inches(0.6), Inches(3.22), W - Inches(1), Inches(0.4),
             font_size=Pt(13), color=MUTED, align=PP_ALIGN.LEFT)

    # meta grid boxes
    meta = [
        ("Report Date", "2026-03-08"),
        ("Authority", "Alaa Atia (AA)"),
        ("Version", "CANONICAL_v1.0"),
        ("Classification", "INTERNAL"),
        ("Analysis Period", "Last 60 Days"),
        ("Systems Covered", "PMS · AlCatara · BIG G · CARSHUNTER"),
    ]
    box_w = Inches(2.0)
    box_h = Inches(0.75)
    cols = 3
    for i, (label, val) in enumerate(meta):
        col = i % cols
        row = i // cols
        lft = Inches(0.6) + col * Inches(2.2)
        tp  = Inches(4.1) + row * Inches(0.9)
        r = add_rect(slide, lft, tp, box_w, box_h, fill_color=MID, line_color=GOLD)
        r.line.width = Pt(0.5)
        add_text(slide, label.upper(), lft + Inches(0.1), tp + Inches(0.05), box_w - Inches(0.2), Inches(0.28),
                 font_size=Pt(7), bold=True, color=MUTED)
        add_text(slide, val, lft + Inches(0.1), tp + Inches(0.32), box_w - Inches(0.2), Inches(0.38),
                 font_size=Pt(11), bold=True, color=WHITE)
    footer_bar(slide)

def s2_exec_summary(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Executive Summary")
    heading(slide, "Governance Architecture Milestone")
    divider(slide)

    paras = [
        "BPB Solutions LTD has achieved a governance architecture milestone that positions the organization at the intersection of institutional-grade decision systems and market-facing operational platforms.",
        "This report documents governance models, architectural innovations, and strategic positioning developed across four interconnected systems: PMS_GOVERN, AlCatara Decision Logic, THE BIG G Doctrine, and CARSHUNTER Platform.",
        "KEY FINDING: BPB has transitioned from authority creation to authority maintenance — a critical phase shift indicating governance maturity and operational readiness."
    ]
    colors = [MUTED, MUTED, WHITE]
    top = Inches(1.6)
    for text, col in zip(paras, colors):
        add_text(slide, text, Inches(0.6), top, W - Inches(1.2), Inches(0.85),
                 font_size=Pt(13), color=col)
        top += Inches(1.02)

    # badges
    badges = [("GOVERNANCE MATURE", GREEN), ("OPERATIONALLY READY", BLUE), ("AUTHORITY MAINTENANCE", GOLD)]
    lft = Inches(0.6)
    for label, col in badges:
        r = add_rect(slide, lft, Inches(4.9), Inches(2.0), Inches(0.38), fill_color=col)
        add_text(slide, label, lft + Inches(0.08), Inches(4.92), Inches(1.85), Inches(0.3),
                 font_size=Pt(9), bold=True,
                 color=DEEP if col == GOLD else WHITE, align=PP_ALIGN.CENTER)
        lft += Inches(2.15)
    footer_bar(slide)

def s3_arch_stack(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 1 · Governance Architecture Stack")
    heading(slide, "The Four-Layer Stack")
    divider(slide)

    layers = [
        ("THE BIG G / G CONCEPT — Governance Doctrine", RGBColor(0x2a,0x1a,0x3e), GOLD),
        ("AlCatara — Decision Logic Model", RGBColor(0x3a,0x2a,0x10), GOLD),
        ("PMS_GOVERN — Execution Governance Engine", MID, WHITE),
        ("Immutable Ledger — Institutional Memory", RGBColor(0x20,0x20,0x35), MUTED),
    ]
    top = Inches(1.65)
    lyr_h = Inches(0.72)
    arrow_h = Inches(0.3)
    for i, (text, bg, fg) in enumerate(layers):
        r = add_rect(slide, Inches(2), top, Inches(9.33), lyr_h, fill_color=bg, line_color=GOLD)
        r.line.width = Pt(0.75)
        add_text(slide, text, Inches(2.1), top + Inches(0.15), Inches(9.1), Inches(0.45),
                 font_size=Pt(14), bold=True, color=fg, align=PP_ALIGN.CENTER)
        top += lyr_h
        if i < 3:
            add_text(slide, "↓", Inches(6.4), top, Inches(0.5), arrow_h,
                     font_size=Pt(18), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            top += arrow_h

    add_text(slide, "Governance before organization.  Authority before role.  Evidence before action.",
             Inches(0.6), Inches(6.3), W - Inches(1.2), Inches(0.35),
             font_size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)
    footer_bar(slide)

def s4_bigg(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Layer 1 · THE BIG G")
    heading(slide, "📐  Governance Doctrine Layer")
    divider(slide)

    principles = [
        ("Authority > Role", "Authority-first governance — authority determines permission, not vice versa."),
        ("Governance First", "Governance infrastructure built before full organizational structure."),
        ("Explicit Commands", "All actions are structured commands with identity verification."),
        ("Domain Separation", "Strict separation prevents cross-domain authority leakage."),
    ]
    cols = 2
    card_w = Inches(5.8)
    card_h = Inches(1.3)
    for i, (title, desc) in enumerate(principles):
        col = i % cols
        row = i // cols
        lft = Inches(0.55) + col * Inches(6.2)
        tp  = Inches(1.65) + row * Inches(1.5)
        r = add_rect(slide, lft, tp, card_w, card_h, fill_color=MID, line_color=GOLD)
        r.line.width = Pt(0.5)
        add_text(slide, title, lft + Inches(0.15), tp + Inches(0.1), card_w - Inches(0.3), Inches(0.38),
                 font_size=Pt(13), bold=True, color=GOLD)
        add_text(slide, desc, lft + Inches(0.15), tp + Inches(0.5), card_w - Inches(0.3), Inches(0.65),
                 font_size=Pt(11), color=MUTED)

    # highlight box
    r = add_rect(slide, Inches(0.55), Inches(4.75), W - Inches(1.1), Inches(0.95),
                 fill_color=RGBColor(0x18,0x14,0x08), line_color=GOLD)
    r.line.width = Pt(1.5)
    add_text(slide, "Strategic Observation: Most organizations build governance reactively. BPB designed governance infrastructure first — positioning it as a platform layer, not an afterthought.",
             Inches(0.75), Inches(4.82), W - Inches(1.5), Inches(0.82),
             font_size=Pt(11), color=WHITE)
    footer_bar(slide)

def s5_alcatara(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Layer 2 · AlCatara")
    heading(slide, "⚡  AlCatara Decision Model")
    divider(slide)

    add_text(slide, "Translates governance doctrine into executable decision structure. Innovation: complexity-aware decision scaling — authority requirements adapt dynamically to risk level.",
             Inches(0.6), Inches(1.6), W - Inches(1.2), Inches(0.6), font_size=Pt(12), color=MUTED)

    rows = [
        ("Energy (E)", "Decision intent", "Command intent"),
        ("State (C)", "Context and conditions", "System state / payload"),
        ("Governance (G)", "Authority validation", "QGED gate"),
        ("Law / Policy (M)", "Rules enforcement", "Signing policy"),
        ("Output (S0–S5)", "Decision state", "ALLOW / BLOCK / HOLD"),
    ]
    headers = ["Component", "Description", "PMS Implementation"]
    col_widths = [Inches(2.8), Inches(3.8), Inches(5.4)]
    top = Inches(2.3)
    row_h = Inches(0.52)
    lft = Inches(0.55)

    # header row
    x = lft
    for h, cw in zip(headers, col_widths):
        add_rect(slide, x, top, cw, row_h, fill_color=RGBColor(0x20,0x18,0x08), line_color=GOLD)
        add_text(slide, h, x + Inches(0.1), top + Inches(0.1), cw - Inches(0.2), row_h - Inches(0.1),
                 font_size=Pt(9), bold=True, color=GOLD, align=PP_ALIGN.LEFT)
        x += cw
    top += row_h

    for r_i, row in enumerate(rows):
        x = lft
        bg = RGBColor(0x16,0x16,0x28) if r_i % 2 == 0 else RGBColor(0x1c,0x1c,0x30)
        for val, cw in zip(row, col_widths):
            add_rect(slide, x, top, cw, row_h, fill_color=bg, line_color=RGBColor(0x33,0x33,0x50))
            col_c = GOLD if val == row[0] else MUTED
            add_text(slide, val, x + Inches(0.1), top + Inches(0.1), cw - Inches(0.2), row_h - Inches(0.12),
                     font_size=Pt(11), bold=(val == row[0]), color=col_c)
            x += cw
        top += row_h
    footer_bar(slide)

def s6_pms(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Layer 3 · PMS_GOVERN")
    heading(slide, "⚙️  Execution Governance Engine")
    divider(slide)

    add_text(slide, "Command → Complexity Score → Authority Gate (QGED) → Execution → Ledger",
             Inches(0.6), Inches(1.62), W - Inches(1.2), Inches(0.48),
             font_size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # two comparison boxes
    for i, (title, body, border) in enumerate([
        ("Traditional RBAC", "User → Role → Permission → Action", MUTED),
        ("PMS Authority Model", "Command → Complexity → Authority → Evidence", GOLD),
    ]):
        lft = Inches(0.55) + i * Inches(6.2)
        r = add_rect(slide, lft, Inches(2.22), Inches(5.9), Inches(1.1),
                     fill_color=MID, line_color=border)
        r.line.width = Pt(1.2)
        add_text(slide, title, lft + Inches(0.15), Inches(2.28), Inches(5.6), Inches(0.38),
                 font_size=Pt(13), bold=True, color=border)
        add_text(slide, body, lft + Inches(0.15), Inches(2.7), Inches(5.6), Inches(0.55),
                 font_size=Pt(11), color=MUTED)

    # highlight
    r = add_rect(slide, Inches(0.55), Inches(3.52), W - Inches(1.1), Inches(1.35),
                 fill_color=RGBColor(0x18,0x14,0x08), line_color=GOLD)
    r.line.width = Pt(1.5)
    add_text(slide, "Governance Gradient:", Inches(0.75), Inches(3.6), Inches(4), Inches(0.36),
             font_size=Pt(12), bold=True, color=GOLD)
    add_text(slide,
             "Complexity 1–3: Delegated execution  |  Complexity 4–7: Dual-signature required  |  Complexity 8–10: AA authority mandatory",
             Inches(0.75), Inches(3.95), W - Inches(1.5), Inches(0.45),
             font_size=Pt(11), color=WHITE)
    add_text(slide, "Compliance: SHA-256 hash-chained ledger  ·  UTC timestamps  ·  Append-only  ·  Emergency override with full audit trail",
             Inches(0.75), Inches(4.4), W - Inches(1.5), Inches(0.38),
             font_size=Pt(10), color=MUTED)
    footer_bar(slide)

def s7_ledger(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Layer 4 · Institutional Memory")
    heading(slide, "🔗  Immutable Ledger")
    divider(slide)

    add_text(slide, "Permanent, auditable record of all governance decisions. Blockchain-style hash-chained. NO entry may ever be modified.",
             Inches(0.6), Inches(1.62), W - Inches(1.2), Inches(0.52),
             font_size=Pt(12), color=MUTED)

    # 3 ledger blocks
    blocks = [
        ("Entry: vehicle_release", "complexity: 7\nsector: trade\nsigner: AA\nstatus: ALLOW", "hash: a3f7b9e2…"),
        ("Entry: payment_approval", "complexity: 9\nsector: finance\nsigner: AA\nprev: a3f7b9e2…", "hash: 9e2c41f7…"),
        ("Entry: contract_mod", "complexity: 10\nsector: system\nsigner: AA\nprev: 9e2c41f7…", "hash: c8b3d5a1…"),
    ]
    blk_w = Inches(3.5)
    blk_h = Inches(1.95)
    top = Inches(2.28)
    x = Inches(0.55)
    for i, (title, body, hsh) in enumerate(blocks):
        r = add_rect(slide, x, top, blk_w, blk_h, fill_color=MID,
                     line_color=RGBColor(0x44,0x44,0x66))
        add_text(slide, title, x + Inches(0.1), top + Inches(0.1), blk_w - Inches(0.2), Inches(0.35),
                 font_size=Pt(10), bold=True, color=GOLD, font_name="Courier New")
        add_text(slide, body, x + Inches(0.1), top + Inches(0.48), blk_w - Inches(0.2), Inches(0.9),
                 font_size=Pt(9), color=WHITE, font_name="Courier New")
        add_text(slide, hsh, x + Inches(0.1), top + Inches(1.55), blk_w - Inches(0.2), Inches(0.3),
                 font_size=Pt(8), color=GOLD, font_name="Courier New")
        x += blk_w
        if i < 2:
            add_text(slide, "→", x, top + Inches(0.7), Inches(0.55), Inches(0.6),
                     font_size=Pt(22), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
            x += Inches(0.55)

    r = add_rect(slide, Inches(0.55), Inches(4.4), W - Inches(1.1), Inches(0.88),
                 fill_color=RGBColor(0x18,0x14,0x08), line_color=GOLD)
    r.line.width = Pt(1.2)
    add_text(slide, "Dispute Resolution: In cross-border automotive trade, the ledger provides objective proof of the authorization chain — critical for resolving conflicts between buyers, sellers, shipping companies, and brokers.",
             Inches(0.75), Inches(4.46), W - Inches(1.5), Inches(0.75), font_size=Pt(10), color=WHITE)
    footer_bar(slide)

def s8_competitive(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 2 · Market Positioning")
    heading(slide, "📊  Competitive Landscape")
    divider(slide)

    rows = [
        ("Governance",      "RBAC",                   "Authority-first command",        "Dynamic risk adaptation"),
        ("Decision Logic",  "Binary allow/deny",       "Complexity-scaled authority",    "Institutional-grade control"),
        ("Audit Trail",     "DB logs (mutable)",       "Hash-chained immutable ledger",  "Regulator-defensible"),
        ("AI Integration",  "Autonomous execution",    "Proposal-only (human gate)",     "Safety-by-design"),
        ("Platform Type",   "Marketplace / listing",   "Governance platform + market",   "Filtered qualified audience"),
    ]
    headers = ["Category", "Traditional", "BPB Architecture", "Advantage"]
    col_widths = [Inches(2.0), Inches(2.8), Inches(4.2), Inches(3.6)]
    top = Inches(1.65)
    row_h = Inches(0.52)
    lft = Inches(0.35)

    x = lft
    for h, cw in zip(headers, col_widths):
        add_rect(slide, x, top, cw, row_h, fill_color=RGBColor(0x20,0x18,0x08), line_color=GOLD)
        add_text(slide, h, x + Inches(0.08), top + Inches(0.1), cw - Inches(0.16), row_h,
                 font_size=Pt(9), bold=True, color=GOLD)
        x += cw
    top += row_h

    for r_i, row in enumerate(rows):
        x = lft
        bg = RGBColor(0x16,0x16,0x28) if r_i % 2 == 0 else RGBColor(0x1c,0x1c,0x30)
        for val, cw in zip(row, col_widths):
            add_rect(slide, x, top, cw, row_h, fill_color=bg, line_color=RGBColor(0x33,0x33,0x50))
            col_c = GOLD if val == row[2] else (GREEN if val == row[3] else MUTED)
            add_text(slide, val, x + Inches(0.08), top + Inches(0.08), cw - Inches(0.16), row_h,
                     font_size=Pt(10), bold=(val in (row[2], row[3])), color=col_c)
            x += cw
        top += row_h
    footer_bar(slide)

def s9_markets(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 2 · Addressable Markets")
    heading(slide, "🌐  Six Target Verticals")
    divider(slide)

    markets = [
        ("🚗  Automotive Trade",    "Provable authorization chain for high-value cross-border transactions."),
        ("📦  Logistics",           "Command governance for multi-party shipment authorization."),
        ("💼  Investment",          "Multi-signature governance for capital deployment decisions."),
        ("🏢  Procurement",         "Policy-driven governance for enterprise purchasing with authority thresholds."),
        ("🤖  AI Control",          "YAI safety architecture — propose-only AI with human authority gate."),
        ("🏥  Regulated Industries","Immutable evidence ledger for financial services, healthcare, compliance."),
    ]
    card_w = Inches(3.9)
    card_h = Inches(1.25)
    cols = 3
    for i, (title, desc) in enumerate(markets):
        col = i % cols
        row = i // cols
        lft = Inches(0.45) + col * Inches(4.15)
        tp  = Inches(1.7) + row * Inches(1.45)
        r = add_rect(slide, lft, tp, card_w, card_h, fill_color=MID, line_color=GOLD)
        r.line.width = Pt(0.5)
        add_text(slide, title, lft + Inches(0.12), tp + Inches(0.1), card_w - Inches(0.24), Inches(0.38),
                 font_size=Pt(11), bold=True, color=GOLD)
        add_text(slide, desc, lft + Inches(0.12), tp + Inches(0.5), card_w - Inches(0.24), Inches(0.68),
                 font_size=Pt(10), color=MUTED)
    footer_bar(slide)

def s10_portfolios(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 3 · Canonical Portfolios")
    heading(slide, "📁  BPB Asset Valuation Tiers")
    divider(slide)

    assets = [
        ("PMS_GOVERN Engine",       "Tier 1 — Core IP",        GOLD, "Exportable · Multi-industry · Regulator-grade · First-mover"),
        ("AlCatara Decision Logic",  "Tier 1 — Core IP",        GOLD, "Complexity-aware · Innovation · Cross-domain"),
        ("THE BIG G Doctrine",       "Tier 1 — Strategic",      GOLD, "Philosophical foundation · Authority sovereignty"),
        ("CARSHUNTER Platform",      "Tier 2 — Revenue",        BLUE, "Authority brand · Filtered audience · Governed ops"),
        ("YAI Safety Layer",         "Tier 1 — Regulatory",     GOLD, "Proposal-only AI · Human gate · Safety-by-design"),
        ("Documentation Archive",    "Tier 3 — Evidence",       MUTED,"Governance record · Audit trail · IP authorship proof"),
    ]
    card_w = Inches(3.9)
    card_h = Inches(1.3)
    cols = 3
    for i, (title, tier, tier_col, desc) in enumerate(assets):
        col = i % cols
        row = i // cols
        lft = Inches(0.45) + col * Inches(4.15)
        tp  = Inches(1.68) + row * Inches(1.5)
        r = add_rect(slide, lft, tp, card_w, card_h, fill_color=MID, line_color=tier_col)
        r.line.width = Pt(0.5)
        add_text(slide, title, lft + Inches(0.12), tp + Inches(0.08), card_w - Inches(0.24), Inches(0.32),
                 font_size=Pt(11), bold=True, color=WHITE)
        add_text(slide, tier, lft + Inches(0.12), tp + Inches(0.4), card_w - Inches(0.24), Inches(0.3),
                 font_size=Pt(11), bold=True, color=tier_col)
        add_text(slide, desc, lft + Inches(0.12), tp + Inches(0.72), card_w - Inches(0.24), Inches(0.5),
                 font_size=Pt(9), color=MUTED)
    footer_bar(slide)

def s11_timeline(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 4 · Operations — Last 60 Days")
    heading(slide, "📅  Governance System Evolution")
    divider(slide)

    milestones = [
        ("2026-01-10", "PMS_GOVERN doctrine formalized",         "Authority chain defined · QGED gate implemented"),
        ("2026-01-22", "AlCatara integrated with PMS",           "Complexity-aware governance activated"),
        ("2026-02-05", "Immutable ledger deployed",              "Hash-chained audit trail operational"),
        ("2026-02-18", "YAI governance framework locked",        "AI safety layer prevents autonomous execution"),
        ("2026-03-04", "CARSHUNTER Authority Upgrade v1.0",      "Public interface reaches 100% authority positioning"),
        ("2026-03-06", "PMS Canonical v1.0 sealed",              "Governance system enters maintenance phase"),
    ]
    top = Inches(1.68)
    item_h = Inches(0.78)
    for date, title, sub in milestones:
        # date label
        add_text(slide, date, Inches(0.5), top + Inches(0.05), Inches(1.3), Inches(0.35),
                 font_size=Pt(9), bold=True, color=GOLD, font_name="Courier New")
        # dot
        add_rect(slide, Inches(1.9), top + Inches(0.16), Inches(0.12), Inches(0.12), fill_color=GOLD)
        # connector line (if not last)
        if date != milestones[-1][0]:
            add_rect(slide, Inches(1.945), top + Inches(0.28), Inches(0.03), item_h - Inches(0.2),
                     fill_color=RGBColor(0x44,0x44,0x44))
        # body box
        r = add_rect(slide, Inches(2.15), top, W - Inches(2.7), item_h - Inches(0.06),
                     fill_color=MID, line_color=RGBColor(0x33,0x33,0x55))
        r.line.width = Pt(0.5)
        add_text(slide, title, Inches(2.3), top + Inches(0.05), W - Inches(3.0), Inches(0.32),
                 font_size=Pt(11), bold=True, color=WHITE)
        add_text(slide, sub, Inches(2.3), top + Inches(0.38), W - Inches(3.0), Inches(0.3),
                 font_size=Pt(9.5), color=MUTED)
        top += item_h
    footer_bar(slide)

def s12_ops(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 4 · High-Risk Operations")
    heading(slide, "🔐  54 Operations · Zero Failed Validations")
    divider(slide)

    rows = [
        ("Vehicle Release Authorization", "23", "6–8",  "Dual-sign / AA"),
        ("Payment Approvals",             "14", "7–9",  "AA Required"),
        ("Contract Modifications",         "8", "8–10", "AA Required"),
        ("Governance System Changes",       "6", "10",  "AA Exclusive"),
        ("Partner Authority Delegation",   "3", "9",    "AA Required"),
    ]
    headers = ["Operation Type", "Count", "Complexity", "Authority Level"]
    col_widths = [Inches(5.4), Inches(1.2), Inches(1.8), Inches(3.6)]
    top = Inches(1.68)
    row_h = Inches(0.5)
    lft = Inches(0.55)

    x = lft
    for h, cw in zip(headers, col_widths):
        add_rect(slide, x, top, cw, row_h, fill_color=RGBColor(0x20,0x18,0x08), line_color=GOLD)
        add_text(slide, h, x + Inches(0.1), top + Inches(0.1), cw - Inches(0.2), row_h,
                 font_size=Pt(9), bold=True, color=GOLD)
        x += cw
    top += row_h

    for r_i, row in enumerate(rows):
        x = lft
        bg = RGBColor(0x16,0x16,0x28) if r_i % 2 == 0 else RGBColor(0x1c,0x1c,0x30)
        for val, cw in zip(row, col_widths):
            add_rect(slide, x, top, cw, row_h, fill_color=bg, line_color=RGBColor(0x33,0x33,0x50))
            add_text(slide, val, x + Inches(0.1), top + Inches(0.08), cw - Inches(0.2), row_h,
                     font_size=Pt(11), color=WHITE if val == row[0] else GOLD)
            x += cw
        top += row_h

    r = add_rect(slide, Inches(0.55), Inches(4.6), W - Inches(1.1), Inches(0.88),
                 fill_color=RGBColor(0x18,0x14,0x08), line_color=GOLD)
    r.line.width = Pt(1.2)
    add_text(slide, "Compliance Record: 12 canonical governance documents  ·  37 BMW vehicle offers processed  ·  8 complete AI-assisted governance sessions  ·  142 WhatsApp records archived",
             Inches(0.75), Inches(4.66), W - Inches(1.5), Inches(0.75), font_size=Pt(10), color=WHITE)
    footer_bar(slide)

def s13_strategy(prs):
    slide = blank_slide(prs)
    gold_bar(slide)
    slide_label(slide, "Part 5 · Strategic Recommendations")
    heading(slide, "🚀  Market Entry Strategy")
    divider(slide)

    phases = [
        ("Phase 1 · Q2 2026", "CARSHUNTER as proof-of-concept for governance-backed platform",     "10 qualified dealer inquiries"),
        ("Phase 2 · Q3 2026", "Package PMS as Governance-as-a-Service for automotive trade networks", "2 pilot implementations"),
        ("Phase 3 · Q4 2026", "Expand to logistics and procurement verticals",                        "5 enterprise clients"),
        ("Phase 4 · Q1 2027", "Position as AI governance infrastructure for regulated industries",    "Regulatory recognition"),
    ]
    card_w = Inches(2.9)
    card_h = Inches(1.6)
    for i, (phase, action, metric) in enumerate(phases):
        lft = Inches(0.45) + i * Inches(3.1)
        r = add_rect(slide, lft, Inches(1.68), card_w, card_h, fill_color=MID, line_color=GOLD)
        r.line.width = Pt(0.8)
        add_text(slide, phase.upper(), lft + Inches(0.12), Inches(1.76), card_w - Inches(0.24), Inches(0.3),
                 font_size=Pt(8), bold=True, color=GOLD)
        add_text(slide, action, lft + Inches(0.12), Inches(2.1), card_w - Inches(0.24), Inches(0.8),
                 font_size=Pt(10), bold=True, color=WHITE)
        add_text(slide, f"Metric: {metric}", lft + Inches(0.12), Inches(2.95), card_w - Inches(0.24), Inches(0.28),
                 font_size=Pt(9), color=MUTED)

    r = add_rect(slide, Inches(0.55), Inches(3.5), W - Inches(1.1), Inches(1.1),
                 fill_color=RGBColor(0x18,0x14,0x08), line_color=GOLD)
    r.line.width = Pt(1.2)
    add_text(slide, "IP Protection: File provisional patent for authority-first architecture  ·  Trademark PMS_GOVERN & AlCatara  ·  Copyright all canonical docs with hash-stamps  ·  Establish trade secrets for QGED gate logic",
             Inches(0.75), Inches(3.57), W - Inches(1.5), Inches(0.95), font_size=Pt(10), color=WHITE)
    footer_bar(slide)

def s14_closing(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, W, Inches(0.12), fill_color=GOLD)
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), fill_color=GOLD)

    add_text(slide, "🏛", Inches(5.8), Inches(0.5), Inches(1.73), Inches(1.1),
             font_size=Pt(48), align=PP_ALIGN.CENTER)

    add_text(slide, "INTERNAL — CONFIDENTIAL  ·  DISTRIBUTION: AA ONLY",
             Inches(0.5), Inches(1.55), W - Inches(1), Inches(0.35),
             font_size=Pt(8), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "BPB Solutions LTD", Inches(0.5), Inches(2.0), W - Inches(1), Inches(1.0),
             font_size=Pt(46), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(5.6), Inches(3.1), Inches(2.13), Inches(0.055), fill_color=GOLD)

    bullets = [
        ("Governance Maturity:", "Institutional-grade architecture — rare for an SME. Closer to financial clearing houses than typical business software."),
        ("Competitive Moat:",     "First-mover in complexity-aware governance scaling. RBAC cannot adapt authority requirements dynamically to risk."),
        ("Current Phase:",        "Authority Maintenance — operational, stable, and ready for qualified counterpart engagement."),
    ]
    top = Inches(3.28)
    for label, body in bullets:
        add_text(slide, label, Inches(0.8), top, Inches(2.3), Inches(0.36),
                 font_size=Pt(11), bold=True, color=GOLD)
        add_text(slide, body, Inches(3.2), top, W - Inches(3.7), Inches(0.42),
                 font_size=Pt(10), color=MUTED)
        top += Inches(0.55)

    add_text(slide, "Hash: SHA256:8f7a3e9b2c1d4f5a…  ·  YAI Governance Analysis System  ·  Supervisor: Alaa Atia (AA)",
             Inches(0.5), Inches(5.4), W - Inches(1), Inches(0.3),
             font_size=Pt(8), color=RGBColor(0x44,0x44,0x66), align=PP_ALIGN.CENTER)

    add_text(slide, "CANONICAL_v1.0  ·  8 MARCH 2026",
             Inches(0.5), H - Inches(0.55), W - Inches(1), Inches(0.35),
             font_size=Pt(9), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── Main ──────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()
    s1_title(prs)
    s2_exec_summary(prs)
    s3_arch_stack(prs)
    s4_bigg(prs)
    s5_alcatara(prs)
    s6_pms(prs)
    s7_ledger(prs)
    s8_competitive(prs)
    s9_markets(prs)
    s10_portfolios(prs)
    s11_timeline(prs)
    s12_ops(prs)
    s13_strategy(prs)
    s14_closing(prs)

    out = "/Users/alaaatia/1. Project AI Manifest/BPB_Governance_Architecture_Report.pptx"
    prs.save(out)
    print(f"✅  Saved: {out}")
    print(f"   Slides: {len(prs.slides)}")
